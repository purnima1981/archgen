"""
PPTX Exporter — converts Diagram JSON → PowerPoint slide

Uses the same output from diagram_builder.build_diagram() so the
canvas layout and PPTX export are always in sync.

Usage:
    from pptx_exporter import export_pptx
    export_pptx(diagram_json, edges, output_path, title=..., subtitle=...)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ═══════════════════════════════════════════════════════════
# COLOR SYSTEM v2 — SEMANTIC
# ═══════════════════════════════════════════════════════════
# Blue/Navy  = TRUST (security, identity, GCP)
# Green      = DATA IN MOTION (pipeline, healthy)
# Orange     = CONTROL (orchestration) + ALERTING
# Purple     = INTELLIGENCE (serving, ML, BI)
# Gray/Slate = EXTERNAL (third-party, outside)
# Teal/Cyan  = OBSERVATION + OUTPUT (monitoring, consumers)

C = {
    # Zone borders
    "gcp":       "1A73E8",
    "security":  "1E3A5F",
    "pipeline":  "137333",
    "medallion": "1B5E20",
    "serving":   "6A1B9A",
    "orch":      "E65100",
    "gcpObs":    "00695C",
    "source":    "546E7A",
    "extId":     "37474F",
    "consumer":  "00838F",
    "extLog":    "455A64",
    "extAlert":  "BF360C",

    # Medallion node borders
    "bronze":    "795548",
    "silver":    "78909C",
    "gold":      "FF8F00",

    # Pipeline stages
    "ingestion": "137333",
    "landing":   "2E7D32",
    "processing":"388E3C",

    # Edges
    "eData":     "2E7D32",
    "eIdentity": "1E3A5F",
    "eControl":  "E65100",
    "eObsInt":   "00695C",
    "eObsExt":   "455A64",
    "eAlert":    "BF360C",
    "eVPN":      "546E7A",

    # UI
    "navy":      "1E3A5F",
    "white":     "FFFFFF",
    "bg":        "F0F4F8",
    "text":      "1E293B",
    "muted":     "64748B",
    "lightBg":   "F8FAFC",
}


def _rgb(hex_str: str) -> RGBColor:
    """Convert 6-char hex to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_rounded_rect(slide, left, top, width, height, fill_hex=None,
                       line_hex=None, line_width=Pt(1.5), line_dash=None,
                       transparency=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.rotation = 0
    # Fill
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_hex)
        if transparency is not None:
            # python-pptx doesn't directly support fill transparency easily,
            # so we use a very light color instead
            pass
    else:
        shape.fill.background()
    # Line
    if line_hex:
        shape.line.color.rgb = _rgb(line_hex)
        shape.line.width = line_width
        if line_dash:
            shape.line.dash_style = line_dash
    else:
        shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text,
                   font_size=10, bold=False, color="1E293B",
                   font_name="Calibri", alignment=PP_ALIGN.LEFT,
                   valign=MSO_ANCHOR.TOP):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    p.font.name = font_name
    p.alignment = alignment
    tf.auto_size = None
    return txBox


def _add_line(slide, x1, y1, x2, y2, color_hex, width=Pt(1.5), dash=None):
    """Add a connector line."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = _rgb(color_hex)
    connector.line.width = width
    if dash:
        connector.line.dash_style = dash
    return connector


# ═══════════════════════════════════════════════════════════
# ZONE / NODE / EDGE PLACEMENT HELPERS
# ═══════════════════════════════════════════════════════════

# Mapping from node group → zone color
GROUP_COLORS = {
    "source":    C["source"],
    "ext-id":    C["extId"],
    "gcp-sec":   C["security"],
    "pipeline":  C["pipeline"],
    "medallion": C["medallion"],
    "serving":   C["serving"],
    "orch":      C["orch"],
    "obs":       C["gcpObs"],
    "consumer":  C["consumer"],
    "ext-alert": C["extAlert"],
    "ext-log":   C["extLog"],
}

# Node-specific overrides
NODE_BORDER_OVERRIDES = {
    "bronze":       C["bronze"],
    "silver":       C["silver"],
    "gold":         C["gold"],
    "gcs":          C["landing"],
    "gcs_raw":      C["landing"],
    "bq_staging":   C["landing"],
    "dataform":     C["processing"],
    "dataflow_proc":C["processing"],
    "dataproc":     C["processing"],
}

# Edge type → color
EDGE_TYPE_COLORS = {
    "data":     C["eData"],
    "identity": C["eIdentity"],
    "control":  C["eControl"],
    "obsInt":   C["eObsInt"],
    "obsExt":   C["eObsExt"],
    "alert":    C["eAlert"],
    "observe":  C["eObsExt"],
}


# ═══════════════════════════════════════════════════════════
# COORDINATE SCALING
# ═══════════════════════════════════════════════════════════
# SVG viewBox 0-1520 x 0-1280 → PPTX 13.3" × 7.5" (LAYOUT_WIDE)

SLIDE_W = 13.3
SLIDE_H = 7.5
HEADER_H = 0.7
FOOTER_H = 0.5
CONTENT_TOP = HEADER_H + 0.05
CONTENT_H = SLIDE_H - HEADER_H - FOOTER_H - 0.1

SVG_W = 1520
SVG_H = 1280


def _sx(v):
    """Scale SVG x → PPTX inches."""
    return 0.25 + (v / SVG_W) * (SLIDE_W - 0.5)


def _sy(v):
    """Scale SVG y → PPTX inches."""
    return CONTENT_TOP + (v / SVG_H) * CONTENT_H


def _sw(v):
    """Scale SVG width → PPTX inches."""
    return (v / SVG_W) * (SLIDE_W - 0.5)


def _sh(v):
    """Scale SVG height → PPTX inches."""
    return (v / SVG_H) * CONTENT_H


# Card dimensions in inches
CARD_W = 0.85
CARD_H = 0.58


# ═══════════════════════════════════════════════════════════
# ZONE DEFINITIONS (matching the React mockup)
# ═══════════════════════════════════════════════════════════

ZONES = [
    # (svg_x, svg_y, svg_w, svg_h, label, color, dashed, filled)
    (820, 30, 420, 170, "CONSUMERS (L8)", C["consumer"], True, False),
    (30, 240, 210, 330, "EXTERNAL IDENTITY", C["extId"], True, False),
    (30, 520, 210, 190, "ON-PREM SOURCE (L1)", C["source"], True, False),
    # GCP boundary
    (410, 220, 1070, 770, "GOOGLE CLOUD PLATFORM", C["gcp"], False, True),
    # Sub-zones inside GCP
    (435, 255, 175, 730, "SECURITY (L2)", C["security"], True, False),
    (635, 395, 420, 570, "DATA PIPELINE (L3→L6)", C["pipeline"], True, False),
    (635, 215, 420, 150, "SERVING (L7)", C["serving"], True, False),
    (1075, 495, 195, 190, "ML / AI", C["serving"], True, False),
    (635, 335, 195, 130, "ORCHESTRATION", C["orch"], True, False),
    (1075, 645, 380, 330, "OBSERVABILITY (GCP)", C["gcpObs"], True, False),
    # External boxes below
    (410, 1060, 330, 195, "EXTERNAL LOGGING", C["extLog"], True, False),
    (830, 1060, 430, 195, "EXTERNAL ALERTING", C["extAlert"], True, False),
]

# ═══════════════════════════════════════════════════════════
# NODE POSITIONS (matching React mockup SVG coords)
# ═══════════════════════════════════════════════════════════

# id → (svg_x, svg_y, name, subtitle, group, emoji)
NODE_POSITIONS = {
    "analysts":    (880, 60,   "Analysts",     "BI Users",       "consumer",  "👤"),
    "executives":  (1060, 60,  "Executives",   "C-suite",        "consumer",  "👔"),
    "oracle":      (70, 560,   "Oracle DB",    "On-prem CDC",    "source",    "🔴"),
    "oracle_db":   (70, 560,   "Oracle DB",    "On-prem CDC",    "source",    "🔴"),
    "entra":       (70, 280,   "Entra ID",     "SSO / OIDC",     "ext-id",    "🔷"),
    "entra_id":    (70, 280,   "Entra ID",     "SSO / OIDC",     "ext-id",    "🔷"),
    "cyberark":    (70, 420,   "CyberArk",     "PAM Vault",      "ext-id",    "🔐"),
    "iam":         (460, 270,  "Cloud IAM",    "Identity",       "gcp-sec",   "🛡️"),
    "cloud_iam":   (460, 270,  "Cloud IAM",    "Identity",       "gcp-sec",   "🛡️"),
    "kms":         (460, 410,  "Cloud KMS",    "CMEK",           "gcp-sec",   "🔑"),
    "cloud_kms":   (460, 410,  "Cloud KMS",    "CMEK",           "gcp-sec",   "🔑"),
    "secret":      (460, 550,  "Secret Mgr",   "Credentials",    "gcp-sec",   "🗝️"),
    "secret_manager":(460,550, "Secret Mgr",   "Credentials",    "gcp-sec",   "🗝️"),
    "vpc":         (460, 690,  "VPC",          "Private Net",    "gcp-sec",   "🌐"),
    "vpcsc":       (460, 830,  "VPC-SC",       "Perimeter",      "gcp-sec",   "🛡️"),
    "vpc_sc":      (460, 830,  "VPC-SC",       "Perimeter",      "gcp-sec",   "🛡️"),
    "datastream":  (660, 880,  "Datastream",   "CDC Stream",     "pipeline",  "⚡"),
    "gcs":         (660, 680,  "GCS Raw",      "Landing Zone",   "pipeline",  "📦"),
    "gcs_raw":     (660, 730,  "GCS Raw",      "Landing Zone",   "pipeline",  "📦"),
    "bq_staging":  (660, 580,  "BQ Staging",   "Staging",        "pipeline",  "📦"),
    "dataform":    (660, 580,  "Dataform",     "SQL ELT (dbt)",  "pipeline",  "⚙️"),
    "dataflow_proc":(660, 530, "Dataflow",     "Processing",     "pipeline",  "⚙️"),
    "dataproc":    (660, 530,  "Dataproc",     "Spark",          "pipeline",  "⚙️"),
    "bronze":      (660, 430,  "Bronze",       "Raw / dedup",    "medallion", "🥉"),
    "silver":      (830, 430,  "Silver",       "Cleaned",        "medallion", "🥈"),
    "gold":        (1000, 430, "Gold",         "Curated",        "medallion", "🥇"),
    "looker":      (660, 250,  "Looker",       "Governed BI",    "serving",   "📊"),
    "vertex":      (1100, 530, "Vertex AI",    "ML Platform",    "serving",   "🤖"),
    "vertex_ai":   (1100, 530, "Vertex AI",    "ML Platform",    "serving",   "🤖"),
    "composer":    (660, 370,  "Composer",     "Airflow DAGs",   "orch",      "🎼"),
    "cloud_composer":(660,370, "Composer",     "Airflow DAGs",   "orch",      "🎼"),
    "monitoring":  (1100, 680, "Monitoring",   "Metrics",        "obs",       "📈"),
    "cloud_monitoring":(1100,680,"Monitoring",  "Metrics",        "obs",       "📈"),
    "logging":     (1100, 830, "Logging",      "Centralized",    "obs",       "📋"),
    "cloud_logging":(1100,830, "Logging",      "Centralized",    "obs",       "📋"),
    "audit":       (1280, 830, "Audit Logs",   "Compliance",     "obs",       "📝"),
    "audit_logs":  (1280, 830, "Audit Logs",   "Compliance",     "obs",       "📝"),
    "pagerduty":   (880, 1110, "PagerDuty",    "Incidents",      "ext-alert", "🚨"),
    "pagerduty_inc":(880,1110, "PagerDuty",    "Incidents",      "ext-alert", "🚨"),
    "wiz":         (1100, 1110,"Wiz CSPM",     "Cloud Security", "ext-alert", "🔒"),
    "wiz_cspm":    (1100,1110, "Wiz CSPM",     "Cloud Security", "ext-alert", "🔒"),
    "splunk":      (460, 1110, "Splunk SIEM",  "Log Analytics",  "ext-log",   "📡"),
    "splunk_siem": (460, 1110, "Splunk SIEM",  "Log Analytics",  "ext-log",   "📡"),
    "dynatrace":   (660, 1110, "Dynatrace",    "APM Traces",     "ext-log",   "🔬"),
    "dynatrace_apm":(660,1110, "Dynatrace",    "APM Traces",     "ext-log",   "🔬"),
    # Additional products that might appear
    "pubsub":      (660, 780,  "Pub/Sub",      "Message Bus",    "pipeline",  "📨"),
    "dataflow_ing":(660, 730,  "Dataflow",     "Stream Ingest",  "pipeline",  "⚡"),
    "cloud_functions":(660,680,"Cloud Func",   "API Pull",       "pipeline",  "⚡"),
    "fivetran":    (660, 780,  "Fivetran",     "Managed ELT",    "pipeline",  "⚡"),
    "looker_studio":(830,250, "Looker Studio","Dashboards",     "serving",   "📊"),
    "cloud_run":   (1060, 370, "Cloud Run",    "API Serving",    "serving",   "🚀"),
    "dataplex":    (460, 1110, "Dataplex",     "Governance",     "obs",       "📋"),
    "data_catalog":(660, 1110, "Data Catalog", "Metadata",       "obs",       "📋"),
    "dataplex_dq": (660, 1060, "Dataplex DQ",  "Data Quality",   "obs",       "📋"),
    "cloud_dlp":   (460, 1060, "Cloud DLP",    "PII Detection",  "obs",       "🔒"),
    "data_scientists":(1060,60,"Data Scientists","ML / Notebooks","consumer",  "🔬"),
    "downstream_sys":(1240,60, "Downstream",   "API Consumers",  "consumer",  "🔗"),
    "cloud_scheduler":(660,420,"Scheduler",    "Cron Triggers",  "orch",      "⏰"),
    "cloud_vpn":   (460, 760,  "Cloud VPN",    "IPSec Tunnel",   "gcp-sec",   "🔒"),
    "cloud_armor": (460, 340,  "Cloud Armor",  "WAF / DDoS",     "gcp-sec",   "🛡️"),
}


# ═══════════════════════════════════════════════════════════
# MAIN EXPORT FUNCTION
# ═══════════════════════════════════════════════════════════

def export_pptx(diagram: dict, output_path: str,
                title: str = "Data Pipeline Architecture",
                subtitle: str = "Generated by ArchGen"):
    """
    Export a Diagram JSON (from diagram_builder) to a PPTX file.

    Args:
        diagram: dict with 'nodes', 'edges', 'title', 'subtitle'
        output_path: where to save the .pptx file
        title: header title
        subtitle: header subtitle
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(C["bg"])

    # ── HEADER BAR ──
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(HEADER_H)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = _rgb(C["navy"])
    hdr.line.fill.background()

    _add_text_box(slide, Inches(0.4), Inches(0.1), Inches(8), Inches(0.4),
                  title, font_size=20, bold=True, color=C["white"])
    _add_text_box(slide, Inches(0.4), Inches(0.42), Inches(8), Inches(0.25),
                  subtitle, font_size=9, color="94A3B8")
    _add_text_box(slide, Inches(SLIDE_W - 2), Inches(0.15), Inches(1.6), Inches(0.4),
                  "ArchGen", font_size=13, bold=True, color=C["white"],
                  alignment=PP_ALIGN.RIGHT)

    # ── FOOTER BAR ──
    ftr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(SLIDE_H - FOOTER_H),
        Inches(SLIDE_W), Inches(FOOTER_H)
    )
    ftr.fill.solid()
    ftr.fill.fore_color.rgb = _rgb(C["navy"])
    ftr.line.fill.background()

    node_count = len(diagram.get("nodes", []))
    edge_count = len(diagram.get("edges", []))

    _add_text_box(slide, Inches(0.4), Inches(SLIDE_H - FOOTER_H + 0.1),
                  Inches(6), Inches(0.3),
                  f"Generated by ArchGen — Enterprise Architecture Intelligence",
                  font_size=8, color="94A3B8")
    _add_text_box(slide, Inches(8), Inches(SLIDE_H - FOOTER_H + 0.1),
                  Inches(5), Inches(0.3),
                  f"{node_count} products · {edge_count} connections    CONFIDENTIAL",
                  font_size=8, color="64748B", font_name="Consolas",
                  alignment=PP_ALIGN.RIGHT)

    # ── ZONE BOXES ──
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    for (zx, zy, zw, zh, zlabel, zcolor, zdashed, zfilled) in ZONES:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(_sx(zx)), Inches(_sy(zy)),
            Inches(_sw(zw)), Inches(_sh(zh))
        )
        if zfilled:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(zcolor)
            # Make very transparent — lighten the color
            shape.fill.fore_color.rgb = _rgb(zcolor)
        else:
            shape.fill.background()

        shape.line.color.rgb = _rgb(zcolor)
        shape.line.width = Pt(1.8)
        if zdashed:
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # Zone label
        _add_text_box(
            slide,
            Inches(_sx(zx) + 0.08), Inches(_sy(zy) - 0.05),
            Inches(len(zlabel) * 0.065 + 0.3), Inches(0.2),
            zlabel, font_size=7, bold=True, color=zcolor
        )

    # ── NODES ──
    nodes = diagram.get("nodes", [])
    node_centers = {}  # id → (center_x_inches, center_y_inches)

    for node in nodes:
        nid = node["id"]
        pos = NODE_POSITIONS.get(nid)
        if not pos:
            continue

        svg_x, svg_y, name, sub, grp, emoji = pos
        border_color = NODE_BORDER_OVERRIDES.get(nid, GROUP_COLORS.get(grp, C["muted"]))

        px = _sx(svg_x)
        py = _sy(svg_y)

        # Store center for edge drawing
        node_centers[nid] = (px + CARD_W / 2, py + CARD_H / 2)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px), Inches(py),
            Inches(CARD_W), Inches(CARD_H)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(C["white"])
        card.line.color.rgb = _rgb(border_color)
        card.line.width = Pt(1.2)
        # Shadow
        card.shadow.inherit = False

        # Emoji
        _add_text_box(slide, Inches(px), Inches(py + 0.02),
                      Inches(CARD_W), Inches(0.22),
                      emoji, font_size=14, alignment=PP_ALIGN.CENTER)

        # Name
        _add_text_box(slide, Inches(px), Inches(py + 0.24),
                      Inches(CARD_W), Inches(0.18),
                      name, font_size=7, bold=True, color=C["text"],
                      alignment=PP_ALIGN.CENTER)

        # Subtitle
        _add_text_box(slide, Inches(px), Inches(py + 0.40),
                      Inches(CARD_W), Inches(0.15),
                      sub, font_size=5.5, color=C["muted"],
                      alignment=PP_ALIGN.CENTER)

    # ── EDGES ──
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    edges = diagram.get("edges", [])
    for edge in edges:
        fid = edge.get("from", "")
        tid = edge.get("to", "")
        if fid not in node_centers or tid not in node_centers:
            continue

        etype = edge.get("edgeType", "data")
        color = EDGE_TYPE_COLORS.get(etype, C["eData"])
        dashed = etype in ("control", "observe", "obsInt", "obsExt", "alert")

        fx, fy = node_centers[fid]
        tx, ty = node_centers[tid]

        conn = slide.shapes.add_connector(
            1,  # straight
            Inches(fx), Inches(fy),
            Inches(tx), Inches(ty)
        )
        conn.line.color.rgb = _rgb(color)
        conn.line.width = Pt(2) if etype == "data" else Pt(1.2)
        if dashed:
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # Edge label
        label = edge.get("label", "")
        if label:
            mx = (fx + tx) / 2
            my_val = (fy + ty) / 2
            # Label background
            lbl_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(mx - 0.22), Inches(my_val - 0.07),
                Inches(0.44), Inches(0.14)
            )
            lbl_shape.fill.solid()
            lbl_shape.fill.fore_color.rgb = _rgb(C["white"])
            lbl_shape.line.color.rgb = _rgb(color)
            lbl_shape.line.width = Pt(0.5)

            _add_text_box(
                slide,
                Inches(mx - 0.22), Inches(my_val - 0.07),
                Inches(0.44), Inches(0.14),
                label, font_size=5, bold=True, color=color,
                font_name="Consolas", alignment=PP_ALIGN.CENTER
            )

    # ── VPN LABEL ──
    if "oracle_db" in node_centers or "oracle" in node_centers:
        oid = "oracle_db" if "oracle_db" in node_centers else "oracle"
        _, oy = node_centers[oid]
        vpn_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(_sx(250)), Inches(oy - 0.07),
            Inches(0.6), Inches(0.16)
        )
        vpn_shape.fill.solid()
        vpn_shape.fill.fore_color.rgb = _rgb(C["white"])
        vpn_shape.line.color.rgb = _rgb(C["source"])
        vpn_shape.line.width = Pt(1)

        _add_text_box(
            slide, Inches(_sx(250)), Inches(oy - 0.07),
            Inches(0.6), Inches(0.16),
            "VPN / DIA", font_size=6, bold=True, color=C["source"],
            font_name="Consolas", alignment=PP_ALIGN.CENTER
        )

    # ── LEGEND ──
    legend_items = [
        (C["eData"],     "Data Flow",     False),
        (C["eIdentity"], "Identity",      True),
        (C["eControl"],  "Orchestration", True),
        (C["gcpObs"],    "Observability", True),
        (C["eAlert"],    "Alerting",      True),
        (C["eObsExt"],   "Ext Logging",   True),
    ]
    leg_x = SLIDE_W - 1.8
    leg_y = CONTENT_TOP + 0.1
    for i, (lc, ll, ld) in enumerate(legend_items):
        ly = leg_y + i * 0.2
        line = slide.shapes.add_connector(
            1,
            Inches(leg_x), Inches(ly + 0.06),
            Inches(leg_x + 0.3), Inches(ly + 0.06)
        )
        line.line.color.rgb = _rgb(lc)
        line.line.width = Pt(1.5)
        if ld:
            line.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        _add_text_box(slide, Inches(leg_x + 0.35), Inches(ly - 0.02),
                      Inches(1.2), Inches(0.18),
                      ll, font_size=6.5, bold=True, color=lc)

    # ── SAVE ──
    prs.save(output_path)
    return output_path


# ═══════════════════════════════════════════════════════════
# STANDALONE TEST
# ═══════════════════════════════════════════════════════════

if __name__ == "__main__":
    # Test with a simple diagram
    test_diagram = {
        "title": "Oracle → BigQuery Pipeline",
        "subtitle": "Test export",
        "nodes": [
            {"id": "oracle_db", "name": "Oracle DB", "x": 70, "y": 560},
            {"id": "datastream", "name": "Datastream", "x": 660, "y": 880},
            {"id": "gcs_raw", "name": "GCS Raw", "x": 660, "y": 730},
            {"id": "dataform", "name": "Dataform", "x": 660, "y": 580},
            {"id": "bronze", "name": "Bronze", "x": 660, "y": 430},
            {"id": "silver", "name": "Silver", "x": 830, "y": 430},
            {"id": "gold", "name": "Gold", "x": 1000, "y": 430},
            {"id": "looker", "name": "Looker", "x": 660, "y": 250},
            {"id": "analysts", "name": "Analysts", "x": 880, "y": 60},
            {"id": "cloud_iam", "name": "Cloud IAM", "x": 460, "y": 270},
            {"id": "cloud_kms", "name": "Cloud KMS", "x": 460, "y": 410},
            {"id": "secret_manager", "name": "Secret Mgr", "x": 460, "y": 550},
            {"id": "vpc", "name": "VPC", "x": 460, "y": 690},
            {"id": "vpc_sc", "name": "VPC-SC", "x": 460, "y": 830},
            {"id": "cloud_composer", "name": "Composer", "x": 660, "y": 370},
            {"id": "cloud_monitoring", "name": "Monitoring", "x": 1100, "y": 680},
            {"id": "cloud_logging", "name": "Logging", "x": 1100, "y": 830},
            {"id": "audit_logs", "name": "Audit Logs", "x": 1280, "y": 830},
            {"id": "pagerduty_inc", "name": "PagerDuty", "x": 880, "y": 1110},
            {"id": "wiz_cspm", "name": "Wiz CSPM", "x": 1100, "y": 1110},
            {"id": "splunk_siem", "name": "Splunk", "x": 460, "y": 1110},
        ],
        "edges": [
            {"from": "oracle_db", "to": "datastream", "label": "CDC", "edgeType": "data"},
            {"from": "datastream", "to": "gcs_raw", "label": "Raw", "edgeType": "data"},
            {"from": "gcs_raw", "to": "dataform", "label": "ELT", "edgeType": "data"},
            {"from": "dataform", "to": "bronze", "label": "Ingest", "edgeType": "data"},
            {"from": "bronze", "to": "silver", "label": "Clean", "edgeType": "data"},
            {"from": "silver", "to": "gold", "label": "Curate", "edgeType": "data"},
            {"from": "gold", "to": "looker", "label": "BI", "edgeType": "data"},
            {"from": "looker", "to": "analysts", "label": "Reports", "edgeType": "data"},
            {"from": "cloud_composer", "to": "dataform", "label": "Trigger", "edgeType": "control"},
            {"from": "cloud_logging", "to": "cloud_monitoring", "label": "Metrics", "edgeType": "observe"},
            {"from": "cloud_monitoring", "to": "pagerduty_inc", "label": "Alerts", "edgeType": "alert"},
            {"from": "cloud_logging", "to": "splunk_siem", "label": "Export", "edgeType": "observe"},
        ],
    }

    out = export_pptx(
        test_diagram,
        "/mnt/user-data/outputs/Oracle_BQ_Pipeline_Python.pptx",
        title="Oracle → BigQuery Data Pipeline",
        subtitle="CDC Replication · Medallion Architecture · Enterprise Observability"
    )
    print(f"✅ Saved to {out}")
